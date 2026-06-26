"""Office / 结构化文档生成（Python 库，落盘到沙箱路径）."""

from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import Any

from sensehub.security.sandbox import assert_filesystem


def _ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def _write_docx(path: Path, title: str, body: str) -> None:
    from docx import Document

    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    for block in body.split("\n\n"):
        text = block.strip()
        if text:
            doc.add_paragraph(text)
    _ensure_parent(path)
    doc.save(str(path))


def _write_xlsx(path: Path, headers: list[str], rows: list[list[Any]]) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    if headers:
        ws.append(headers)
    for row in rows:
        ws.append(list(row))
    _ensure_parent(path)
    wb.save(str(path))


def _normalize_slide_texts(slides: list[Any]) -> list[str]:
    texts: list[str] = []
    for item in slides:
        if isinstance(item, dict):
            title = str(item.get("title") or "").strip()
            body = str(item.get("content") or item.get("body") or "").strip()
            if title and body:
                texts.append(f"{title}\n{body}")
            elif title:
                texts.append(title)
            elif body:
                texts.append(body)
            continue
        text = str(item).strip()
        if text:
            texts.append(text)
    return texts


def _write_pptx(path: Path, title: str, slides: list[str]) -> None:
    from pptx import Presentation

    prs = Presentation()
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    for content in slides:
        text = content.strip()
        if not text:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = text.split("\n", 1)[0][:80]
        body = slide.placeholders[1]
        body.text = text
    _ensure_parent(path)
    prs.save(str(path))


def generate_document(params: dict[str, Any]) -> dict[str, Any]:
    """生成 docx/xlsx/pptx/txt/csv/md 并写入沙箱路径."""
    fmt = str(params.get("format") or params.get("type") or "txt").lower().strip().lstrip(".")
    rel_path = str(params.get("path") or params.get("filename") or "").strip()
    if not rel_path:
        raise ValueError("path 不能为空")
    if not rel_path.lower().endswith(f".{fmt}"):
        rel_path = f"{rel_path}.{fmt}"

    path = assert_filesystem(rel_path, "write")
    title = str(params.get("title") or "").strip()
    body = str(params.get("content") or params.get("body") or "").strip()

    if fmt in {"txt", "md", "markdown"}:
        _ensure_parent(path)
        path.write_text(body or title, encoding="utf-8")
    elif fmt == "csv":
        headers = params.get("headers") or []
        rows = params.get("rows") or []
        if not isinstance(headers, list):
            headers = []
        if not isinstance(rows, list):
            rows = []
        buf = io.StringIO()
        writer = csv.writer(buf)
        if headers:
            writer.writerow([str(h) for h in headers])
        for row in rows:
            if isinstance(row, list):
                writer.writerow([str(c) for c in row])
        _ensure_parent(path)
        path.write_text(buf.getvalue(), encoding="utf-8")
    elif fmt == "docx":
        try:
            _write_docx(path, title, body)
        except ImportError as exc:
            raise RuntimeError("未安装 python-docx，请 pip install python-docx") from exc
    elif fmt in {"xlsx", "excel"}:
        headers = params.get("headers") or []
        rows = params.get("rows") or []
        if not isinstance(headers, list):
            headers = []
        if not isinstance(rows, list):
            rows = []
        try:
            _write_xlsx(path, [str(h) for h in headers], rows)
        except ImportError as exc:
            raise RuntimeError("未安装 openpyxl，请 pip install openpyxl") from exc
    elif fmt in {"pptx", "ppt"}:
        slides = params.get("slides")
        if isinstance(slides, list) and slides:
            slide_texts = _normalize_slide_texts(slides)
        else:
            slide_texts = [body] if body else ([title] if title else [" "])
        try:
            _write_pptx(path, title, slide_texts)
        except ImportError as exc:
            raise RuntimeError("未安装 python-pptx，请 pip install python-pptx") from exc
    else:
        raise ValueError(f"不支持的文档格式: {fmt}")

    return {
        "path": str(path),
        "format": fmt,
        "bytes": path.stat().st_size,
        "method": "generate_document",
    }
